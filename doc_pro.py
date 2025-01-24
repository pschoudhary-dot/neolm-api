import os
import re
import asyncio
import logging
from typing import List
from pydantic import BaseModel
import hashlib
from urllib.parse import urlparse
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.prompts import ChatPromptTemplate
from langchain_groq import ChatGroq
from langchain.chains.retrieval import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from PIL import Image
from dotenv import load_dotenv

load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# --- Configuration ---
class Config:
    CHROMA_DIR = "./chroma_db"
    PROCESSED_DIR = "./processed_data"
    SUPPORTED_TYPES = {
        'audio': ['mp3', 'wav', 'ogg'],
        'video': ['mp4', 'mov', 'avi'],
        'image': ['jpg', 'jpeg', 'png'],
        'doc': ['pdf', 'docx', 'txt', 'csv', 'xlsx', 'pptx'],
        'url': ['url']
    }
    EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
    LLM_MODEL = "llama3-70b-8192"
    WHISPER_MODEL = "openai/whisper-large-v3-turbo"
    CHUNK_SIZE = 1000
    CHUNK_OVERLAP = 200


# --- Data Models ---
class DocumentMetadata(BaseModel):
    user_id: str
    doc_type: str
    doc_name: str
    processed_text: str
    sha_hash: str


# --- File Processors ---
class FileProcessor:
    def __init__(self, user_id: str):
        self.user_id = user_id
        os.makedirs(Config.PROCESSED_DIR, exist_ok=True)

    @staticmethod
    def _save_metadata(metadata: DocumentMetadata):
        filename = f"{metadata.sha_hash}.json"
        path = os.path.join(Config.PROCESSED_DIR, filename)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(metadata.model_dump_json())
        logger.info(f"Metadata saved: {filename}")

    def process_file(self, file_path: str) -> DocumentMetadata:
        raise NotImplementedError


class AudioProcessor(FileProcessor):
    def process_file(self, file_path: str) -> DocumentMetadata:
        from transformers import pipeline
        logger.info(f"Processing audio file: {file_path}")

        pipe = pipeline(
            "automatic-speech-recognition",
            model=Config.WHISPER_MODEL,
            device_map="auto"
        )
        text = pipe(file_path)["text"]

        sha_hash = hashlib.sha256(text.encode()).hexdigest()
        return DocumentMetadata(
            user_id=self.user_id,
            doc_type="audio",
            doc_name=os.path.basename(file_path),
            processed_text=text,
            sha_hash=sha_hash
        )


class VideoProcessor(AudioProcessor):
    def process_file(self, file_path: str) -> DocumentMetadata:
        from moviepy.editor import VideoFileClip
        import tempfile

        logger.info(f"Processing video file: {file_path}")
        with tempfile.NamedTemporaryFile(suffix=".wav") as tmpfile:
            video = VideoFileClip(file_path)
            video.audio.write_audiofile(tmpfile.name)
            return super().process_file(tmpfile.name)


class ImageProcessor(FileProcessor):
    def process_file(self, file_path: str) -> DocumentMetadata:
        from transformers import BlipProcessor, BlipForConditionalGeneration

        processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-large")
        model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-large")

        raw_image = Image.open(file_path).convert('RGB')
        inputs = processor(raw_image, return_tensors="pt")
        out = model.generate(**inputs)
        text = processor.decode(out[0], skip_special_tokens=True)

        sha_hash = hashlib.sha256(text.encode()).hexdigest()
        return DocumentMetadata(
            user_id=self.user_id,
            doc_type="image",
            doc_name=os.path.basename(file_path),
            processed_text=text,
            sha_hash=sha_hash
        )


class DocumentProcessor(FileProcessor):
    def process_file(self, file_path: str) -> DocumentMetadata:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""

        if ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = "\n".join([page.extract_text() for page in reader.pages])
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == '.csv':
            import pandas as pd
            df = pd.read_csv(file_path)
            text = df.to_markdown()
        elif ext == '.xlsx':
            import pandas as pd
            df = pd.read_excel(file_path)
            text = df.to_markdown()
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            text = "\n".join(text)
        else:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()

        text = text.encode('utf-8', 'replace').decode('utf-8')
        sha_hash = hashlib.sha256(text.encode()).hexdigest()
        return DocumentMetadata(
            user_id=self.user_id,
            doc_type="document",
            doc_name=os.path.basename(file_path),
            processed_text=text,
            sha_hash=sha_hash
        )


class URLProcessor(FileProcessor):
    async def process_file_async(self, url: str) -> DocumentMetadata:
        from crawl4ai import AsyncWebCrawler

        try:
            async with AsyncWebCrawler(
                    browser_config={
                        "headless": True,
                        "user_agent": "Mozilla/5.0 (compatible; MyCrawler/1.0)",
                        "timeout": 30
                    }
            ) as crawler:
                result = await crawler.arun(
                    url=url,
                    run_config={
                        "cache_mode": "bypass",
                        "extraction_strategy": "auto",
                        "js_code": [
                            "window.scrollTo(0, document.body.scrollHeight)",
                            "new Promise(resolve => setTimeout(resolve, 2000))"
                        ]
                    }
                )

                if not result.cleaned_text:
                    raise ValueError("No meaningful content extracted")

                text = self._clean_content(result.cleaned_text)
                metadata = self._create_metadata(text, url)
                self._save_metadata(metadata)
                return metadata

        except Exception as e:
            logger.error(f"Web crawling error: {str(e)}")
            raise

    def _clean_content(self, text: str) -> str:
        patterns = [
            r'<script\b[^>]*>.*?</script>',
            r'<style\b[^>]*>.*?</style>',
            r'\bAdvertisement\b.*?\n',
            r'\s{2,}',
            r'\[.*?\]'
        ]
        for pattern in patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE | re.DOTALL)
        return ' '.join(text.split()[:2000])

    def _create_metadata(self, text: str, url: str) -> DocumentMetadata:
        parsed = urlparse(url)
        domain = parsed.netloc.replace('www.', '').split('.')[0].title()
        sha_hash = hashlib.sha256(text.encode()).hexdigest()
        return DocumentMetadata(
            user_id=self.user_id,
            doc_type="url",
            doc_name=f"{domain}_content",
            processed_text=text,
            sha_hash=sha_hash
        )

    def process_file(self, url: str) -> DocumentMetadata:
        return asyncio.run(self.process_file_async(url))


# --- Processing Pipeline ---
def process_file(file_path: str, user_id: str) -> DocumentMetadata:
    try:
        parsed = urlparse(file_path)
        if parsed.scheme in ('http', 'https'):
            ext = 'url'
        else:
            ext = os.path.splitext(file_path)[1].lower().lstrip('.')
    except:
        ext = 'url'

    processor_map = {
        'audio': AudioProcessor,
        'video': VideoProcessor,
        'image': ImageProcessor,
        'doc': DocumentProcessor,
        'url': URLProcessor
    }

    for file_type, exts in Config.SUPPORTED_TYPES.items():
        if ext in exts or (file_type == 'url' and ext == 'url'):
            processor = processor_map[file_type](user_id)
            metadata = processor.process_file(file_path)
            processor._save_metadata(metadata)
            return metadata

    raise ValueError(f"Unsupported file type: {ext}")


# --- Vector Store Management ---
class VectorStoreManager:
    def __init__(self):
        self.embedding = HuggingFaceEmbeddings(
            model_name=Config.EMBEDDING_MODEL,
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )
        self.vector_store = Chroma(
            persist_directory=Config.CHROMA_DIR,
            embedding_function=self.embedding
        )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.CHUNK_SIZE,
            chunk_overlap=Config.CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " "]
        )

    def add_documents(self, documents: List[DocumentMetadata]):
        chunks = []
        for doc in documents:
            doc_chunks = self.text_splitter.create_documents(
                [doc.processed_text],
                [doc.model_dump()]
            )
            chunks.extend(doc_chunks)

        self.vector_store.add_documents(chunks)
        logger.info(f"Added {len(chunks)} chunks to vector store")

    def get_retriever(self):
        return self.vector_store.as_retriever(
            search_type="mmr",
            search_kwargs={"k": 5, "score_threshold": 0.7}
        )


# --- Chat Interface ---
class ChatManager:
    def __init__(self):
        self.llm = ChatGroq(
            temperature=0.3,
            model_name=Config.LLM_MODEL,
            api_key=os.getenv("GROQ_API_KEY")
        )
        self.vector_store = VectorStoreManager()

        self.prompt = ChatPromptTemplate.from_messages([
            ("system", """you are very helpful and intelligent assistant help the user using context. Follow these rules:
            1. For authorship questions, check 'doc_name' and text content
            2. If unsure, state "Information not found in documents"
            3. Always cite sources using format: [Source: doc_name]

            Context: {context}

            Question: {input}
            Answer:""")
        ])

        document_chain = create_stuff_documents_chain(self.llm, self.prompt)
        self.retrieval_chain = create_retrieval_chain(
            self.vector_store.get_retriever(),
            document_chain
        )

    def chat(self, question: str):
        try:
            response = self.retrieval_chain.invoke({"input": question})
            answer = response["answer"]
            sources = list({doc.metadata['doc_name'] for doc in response["context"]})
            return f"{answer}\n\nSources: {', '.join(sources)}"
        except Exception as e:
            logger.error(f"Error generating answer: {str(e)}")
            return "Sorry, I encountered an error processing your request."


# --- Main Application ---
def main():
    vector_mgr = VectorStoreManager()
    chat_mgr = ChatManager()
    user_id = "user_123"

    files_to_process = [
        ("./documents/voice.mp3", "audio"),
        ("./documents/LLM_Example.docx", "doc"),
        ("https://langchain-ai.github.io/langgraph/concepts/memory/", "url"),
        ("./documents/LLM.png", "image")
    ]

    processed_docs = []
    for file_path, _ in files_to_process:
        try:
            logger.info(f"Processing {file_path}...")
            metadata = process_file(file_path, user_id)
            processed_docs.append(metadata)
            logger.info(f"Processed {file_path} successfully")
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")

    vector_mgr.add_documents(processed_docs)

    while True:
        question = input("\nAsk a question (q to quit): ")
        if question.lower() == 'q':
            break
        response = chat_mgr.chat(question)
        print(f"\nAssistant: {response}")


if __name__ == "__main__":
    main()